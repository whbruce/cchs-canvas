import wget
import tempfile
import pptx

def is_useful_announcement(self, title):
    if title.startswith("****"):
        return False
    elif title.startswith("Attendance"):
        return False
    return True


def get_check_in_time(self, date):
    courses = [5237, 4843, 5237, 4843, 5237]
    day = date.weekday()
    if (day > 4):
        return None
    course_id = courses[day]
    start_date = date - timedelta(hours=12)
    end_date = date + timedelta(hours=12)
    context_code = "course_" + str(course_id)
    for a in self.canvas.get_announcements(context_codes=[context_code], start_date=start_date.strftime("%Y-%m-%d"), end_date=end_date.strftime("%Y-%m-%d")):
        a.course_id = int(a.context_code[7:])
        course = self.course_short_name(self.course_dict[int(a.context_code[7:])])
        try:
            for e in a.get_topic_entries():
                if (e.user_id == self.user_id):
                    t = convert_date(e.created_at).astimezone(pytz.timezone('US/Pacific'))
                    return t
        except exceptions.Forbidden:
            pass
    return None

def get_announcements(self):
    courses=[]
    for id in self.course_dict:
        courses.append("course_" + str(id))
    today = datetime.today().astimezone(pytz.timezone('US/Pacific'))
    start_date = today - timedelta(hours=12)
    start_date = start_date.strftime("%Y-%m-%d")
    announcements=[]
    for a in self.canvas.get_announcements(context_codes=courses, start_date=start_date):
        if self.is_useful_announcement(a.title):
            course = self.course_short_name(self.course_dict[int(a.context_code[7:])])
            announcements.append(Announcement(course, a.title, a.message, a.posted_at))
    return announcements



def download_file(self, file):
    path = os.path.join(tempfile.gettempdir(), file.filename)
    print(path)
    print(file.url)
    if not os.path.exists(path):
        try:
            wget.download(file.url, out=path)
        except:
            return None
    return path

def get_course(self, name):
    for course in self.courses:
        print(dir(course))
        if name in course.name:
            return course
    return None

def get_english_notes(self):
    date = datetime.today().astimezone(pytz.timezone('US/Pacific'))
    course = self.get_course("English")
    pptx_mime_type =  "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    files = course.get_files(content_types=pptx_mime_type, sort="created_at", order="desc")
    date_format = "%#m.%d.pptx" if "win" in sys.platform else "%-m.%d.pptx"
    expected_filename = date.strftime(date_format)
    print(expected_filename)
    for file in files:
            if file.filename <= expected_filename:
                break
    else:
        return None
    filename = file.filename
    print(filename)
    download_path = filename
    download_path = self.download_file(file)
    if not download_path:
        return None
    print(download_path)
    buffer = io.StringIO()
    print("### English Notes (%s)" % (filename), file=buffer)
    presentation = pptx.Presentation(download_path)
    slides = presentation.slides
    for slide in slides:
        header = True
        for shape in slide.shapes:
            # print(" - %s %s" % (shape, shape.has_text_frame))
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                        if (paragraph.text):
                            marker = "####" if header else "*"
                            print("%s %s" % (marker, paragraph.text), file=buffer)
                            header = False
    markdown_text = buffer.getvalue()
    buffer.close()
    return markdown_text

def run_submission_report(self):
    submissions = {}
    start_date = datetime(2020, 3, 16).astimezone(pytz.timezone('US/Pacific'))
    end_date = datetime.today().astimezone(pytz.timezone('US/Pacific'))
    num_submissions = 0
    for course in self.courses:
        course_name = self.course_short_name(course)
        if course_name:
            assignments = course.get_assignments(order_by="due_at", include=["submission"])
            for a in assignments:
                assignment = Assignment(course_name, a)
                if assignment.is_valid:
                    due_date = assignment.get_due_date()
                    if due_date > start_date and due_date < end_date and assignment.is_submitted():
                        submission_date = assignment.get_submission_date().astimezone(pytz.timezone('US/Pacific'))
                        hour = submission_date.hour
                        if hour < 10:
                            hour = 10
                        if hour in submissions:
                            submissions[hour] = submissions[hour] + 1
                        else:
                            submissions[hour] = 1
                        num_submissions = num_submissions + 1
                        print("%s %s is early" % (course, assignment.get_name()))
                        print("%s %s" % (due_date.strftime("%m/%d"), assignment.get_submission_date().astimezone(pytz.timezone('US/Pacific'))))
    print(len(submissions))
    max_hour = max(k for k, v in submissions.items())
    print(num_submissions)
    cum_pc = 0
    for hour in sorted(submissions.keys()):
        pc = (100 * submissions[hour]) / num_submissions
        cum_pc = cum_pc + pc
        print("%2d: %2d %2d %3d" % (hour, submissions[hour], pc, cum_pc))

